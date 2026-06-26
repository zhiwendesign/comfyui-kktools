"""
ComfyUI OpenMAIC 节点 - 独立版导入课件节点

完全本地化处理，不依赖外部服务：
1. PPTX 导入（使用 LibreOffice/PowerPoint 转 PDF）
2. 本地讲稿匹配（基于规则，不需要 AI）
3. 动作生成（spotlight + laser + speech）

依赖：
- LibreOffice 或 PowerPoint
- PyMuPDF（用于解析 PDF 和生成图片）
- python-pptx（用于提取 PPTX 文本）
"""

import os
import re
import json
import tempfile
import subprocess
import base64
import shutil
from typing import Tuple, Dict, Any, Optional, List
from xml.etree import ElementTree as ET

NODE_CLASS_MAPPINGS = {}
NODE_DISPLAY_NAME_MAPPINGS = {}


def register_node(cls):
    NODE_CLASS_MAPPINGS[cls.__name__] = cls
    NODE_DISPLAY_NAME_MAPPINGS[cls.__name__] = cls.DISPLAY_NAME
    return cls


def clamp(value: float, min_val: float, max_val: float) -> float:
    return max(min_val, min(max_val, value))


# ============== PPTX 解析相关 ==============

def find_libreoffice() -> Optional[str]:
    candidates = ['soffice', 'libreoffice',
                  'C:\\Program Files\\LibreOffice\\program\\soffice.exe',
                  'C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe']
    for candidate in candidates:
        try:
            subprocess.run([candidate, '--version'], capture_output=True, timeout=5)
            return candidate
        except Exception:
            continue
    return None


def find_powerpoint() -> Optional[str]:
    candidates = ['C:\\Program Files\\Microsoft Office\\root\\Office16\\POWERPNT.EXE',
                  'C:\\Program Files (x86)\\Microsoft Office\\root\\Office16\\POWERPNT.EXE']
    for candidate in candidates:
        if os.path.exists(candidate):
            return candidate
    return None


def strip_html(text: str) -> str:
    text = re.sub(r'<style[^>]*>.*?</style>', ' ', text, flags=re.DOTALL)
    text = re.sub(r'<script[^>]*>.*?</script>', ' ', text, flags=re.DOTALL)
    text = re.sub(r'<[^>]+>', ' ', text)
    text = text.replace('&nbsp;', ' ').replace('&amp;', '&')
    text = text.replace('&lt;', '<').replace('&gt;', '>')
    text = text.replace('&quot;', '"').replace('&#39;', "'")
    return re.sub(r'\s+', ' ', text).strip()


def extract_pptx_text(pptx_path: str) -> List[Dict[str, Any]]:
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slides_data = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_info = {'index': slide_num, 'text': '', 'note': '', 'shapes': []}
            texts = []

            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    texts.append(shape.text.strip())
                    try:
                        slide_info['shapes'].append({
                            'text': shape.text.strip(),
                            'left': shape.left / 914400,
                            'top': shape.top / 914400,
                            'width': shape.width / 914400,
                            'height': shape.height / 914400,
                            'slide_width': prs.slide_width / 914400,
                            'slide_height': prs.slide_height / 914400
                        })
                    except Exception:
                        pass

            slide_info['text'] = ' '.join(texts)

            if slide.has_notes_slide:
                notes_text_frame = slide.notes_slide.notes_text_frame
                if notes_text_frame and notes_text_frame.text.strip():
                    slide_info['note'] = notes_text_frame.text.strip()

            slides_data.append(slide_info)
        return slides_data
    except ImportError:
        return []


def extract_pptx_text_fallback(pptx_path: str) -> List[Dict[str, Any]]:
    import zipfile
    slides_data = []
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            slide_files = sorted([f for f in zf.namelist() if re.match(r'ppt/slides/slide\d+\.xml', f)])
            for slide_num, slide_file in enumerate(slide_files, 1):
                slide_info = {'index': slide_num, 'text': '', 'note': '', 'shapes': []}
                with zf.open(slide_file) as f:
                    content = f.read().decode('utf-8', errors='ignore')
                    slide_info['text'] = ' '.join(strip_html(content).split())
                notes_file = f'ppt/notesSlides/notesSlide{slide_num}.xml'
                if notes_file in zf.namelist():
                    with zf.open(notes_file) as f:
                        slide_info['note'] = strip_html(f.read().decode('utf-8', errors='ignore'))
                slides_data.append(slide_info)
    except Exception:
        pass
    return slides_data


def convert_pptx_to_pdf(pptx_path: str, output_dir: str) -> str:
    pdf_path = os.path.join(output_dir, os.path.splitext(os.path.basename(pptx_path))[0] + '.pdf')

    soffice = find_libreoffice()
    if soffice:
        try:
            subprocess.run([soffice, '--headless', '--nologo', '--nofirststartwizard',
                          '--norestore', '--convert-to', 'pdf', '--outdir', output_dir, pptx_path],
                         capture_output=True, timeout=120)
            if os.path.exists(pdf_path):
                return pdf_path
        except Exception:
            pass

    powerpoint = find_powerpoint()
    if powerpoint:
        script = f'''
$ErrorActionPreference = 'Stop'
$pptxPath = "{pptx_path}"
$pdfPath = "{pdf_path}"
$powerPoint = $null
$presentation = $null
try {{
    $powerPoint = New-Object -ComObject PowerPoint.Application
    $powerPoint.DisplayAlerts = 1
    $presentation = $powerPoint.Presentations.Open($pptxPath, $true, $false, $false)
    $presentation.SaveAs($pdfPath, 32)
}}
finally {{
    if ($presentation) {{ $presentation.Close() }}
    if ($powerPoint) {{ $powerPoint.Quit() }}
    [GC]::Collect()
    [GC]::WaitForPendingFinalizers()
}}
'''
        script_path = os.path.join(output_dir, 'convert.ps1')
        with open(script_path, 'w', encoding='utf-8') as f:
            f.write(script)
        try:
            subprocess.run(['powershell.exe', '-NoProfile', '-NonInteractive', '-ExecutionPolicy', 'Bypass',
                          '-File', script_path], capture_output=True, timeout=120)
            if os.path.exists(pdf_path):
                return pdf_path
        except Exception:
            pass

    raise Exception("无法将 PPTX 转换为 PDF。请安装 LibreOffice 或 PowerPoint。")


def parse_pdf_to_images(pdf_path: str, output_dir: str) -> List[str]:
    image_paths = []
    try:
        import fitz
        doc = fitz.open(pdf_path)
        for page_num in range(len(doc)):
            page = doc[page_num]
            mat = fitz.Matrix(2.0, 2.0)
            pix = page.get_pixmap(matrix=mat)
            img_path = os.path.join(output_dir, f'page_{page_num + 1:03d}.png')
            pix.save(img_path)
            image_paths.append(img_path)
        doc.close()
        return image_paths
    except ImportError:
        try:
            from pdf2image import convert_from_path
            images = convert_from_path(pdf_path, dpi=150)
            for i, img in enumerate(images):
                img_path = os.path.join(output_dir, f'page_{i + 1:03d}.png')
                img.save(img_path, 'PNG')
                image_paths.append(img_path)
            return image_paths
        except ImportError:
            raise Exception("请安装 PyMuPDF (pip install pymupdf) 或 pdf2image (pip install pdf2image)")
    return image_paths


def parse_pdf_extract_text(pdf_path: str) -> List[Dict[str, Any]]:
    slides_data = []
    try:
        import fitz
        doc = fitz.open(pdf_path)
        for page_num in range(len(doc)):
            page = doc[page_num]
            text_dict = page.get_text("dict")
            slide_info = {
                'index': page_num + 1,
                'text': '',
                'width': page.rect.width,
                'height': page.rect.height,
                'shapes': []
            }
            texts = []
            for block in text_dict.get("blocks", []):
                if block.get("type") == 0:
                    for line in block.get("lines", []):
                        for span in line.get("spans", []):
                            text = span.get("text", "").strip()
                            if text:
                                texts.append(text)
                                bbox = span.get("bbox", [0, 0, 0, 0])
                                slide_info['shapes'].append({
                                    'text': text,
                                    'x': bbox[0] / page.rect.width * 100,
                                    'y': bbox[1] / page.rect.height * 100,
                                    'w': (bbox[2] - bbox[0]) / page.rect.width * 100,
                                    'h': (bbox[3] - bbox[1]) / page.rect.height * 100
                                })
            slide_info['text'] = ' '.join(texts)
            slides_data.append(slide_info)
        doc.close()
        return slides_data
    except ImportError:
        raise Exception("请安装 PyMuPDF (pip install pymupdf)")
    return slides_data


# ============== 讲稿匹配相关 ==============

def split_speech_segments(text: str, max_chars: int = 260) -> List[str]:
    """将文本分割为语音片段"""
    if not text:
        return []
    text = text.replace('\r\n', '\n').replace('\n\n\n', '\n\n').strip()
    paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
    segments = []
    for para in paragraphs if paragraphs else [text]:
        if len(para) <= max_chars:
            segments.append(para)
        else:
            # 按句子分割
            sentences = re.split(r'(?<=[。！？!?;；.])', para)
            current = ''
            for sent in sentences:
                if len(current) + len(sent) <= max_chars:
                    current += sent
                else:
                    if current:
                        segments.append(current.strip())
                    current = sent
            if current.strip():
                segments.append(current.strip())
    return [s for s in segments if s.strip()]


def match_script_to_slides(deck: Dict, script: str) -> List[Dict]:
    """本地讲稿匹配 - 不依赖 AI"""
    slides = deck.get('slides', [])
    trimmed_script = script.strip()

    if not trimmed_script:
        # 无讲稿，使用每页的 note 或 text
        return [{
            'slideIndex': s.get('index', i + 1),
            'text': s.get('note') or s.get('text', ''),
            'segments': split_speech_segments(s.get('note') or s.get('text', '')),
            'source': 'note' if s.get('note') else ('text' if s.get('text') else 'empty')
        } for i, s in enumerate(slides)]

    # 按页面标记分割（支持：第1页、Slide 1、1/10 等格式）
    marker_pattern = r'^(?:第\s*(\d+)\s*(?:页|張|P)|(?:slide|page|p)\s*(\d+))\s*[:：.-]?\s*$'
    lines = trimmed_script.replace('\r\n', '\n').split('\n')
    buckets = {}
    current_page = None

    for line in lines:
        match = re.match(marker_pattern, line.strip(), re.IGNORECASE)
        if match:
            page_num = int(match.group(1) or match.group(2))
            current_page = page_num
            buckets[current_page] = []
        elif current_page is not None:
            buckets[current_page].append(line)

    if buckets:
        # 成功按标记分割
        return [{
            'slideIndex': i + 1,
            'text': '\n'.join(buckets.get(i + 1, [])),
            'segments': split_speech_segments('\n'.join(buckets.get(i + 1, []))),
            'source': 'marker'
        } for i in range(len(slides))]

    # 按分隔符分割（---、===、***）
    parts = re.split(r'\n\s*(?:---+|===+|\*\*\*+)\s*\n', trimmed_script)
    if len(parts) == len(slides):
        return [{
            'slideIndex': slides[i].get('index', i + 1),
            'text': parts[i].strip(),
            'segments': split_speech_segments(parts[i].strip()),
            'source': 'separator'
        } for i in range(len(slides))]

    # 自动分配：平均分配
    per_slide = trimmed_script
    if len(trimmed_script) > len(slides) * 500:
        chunk_size = len(trimmed_script) // len(slides)
        parts = [trimmed_script[i:i+chunk_size] for i in range(0, len(trimmed_script), chunk_size)]
    else:
        # 按段落数均分
        paras = [p for p in re.split(r'\n\n+', trimmed_script) if p.strip()]
        paras_per_slide = max(1, len(paras) // len(slides))
        parts = ['\n\n'.join(paras[i:i+paras_per_slide]) for i in range(0, len(paras), paras_per_slide)]

    while len(parts) < len(slides):
        parts.append('')
    parts = parts[:len(slides)]

    return [{
        'slideIndex': slides[i].get('index', i + 1),
        'text': parts[i].strip(),
        'segments': split_speech_segments(parts[i].strip()),
        'source': 'auto'
    } for i in range(len(slides))]


# ============== 动作生成相关 ==============

def nanoid(length: int = 8) -> str:
    """生成简单 ID"""
    import random
    chars = 'abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789'
    return ''.join(random.choices(chars, k=length))


def build_imported_lecture_actions(deck: Dict, matches: List[Dict]) -> List[Dict]:
    """生成讲解动作"""
    slides = deck.get('slides', [])
    action_plans = []

    for slide in slides:
        slide_index = slide.get('index', 1)
        match = next((m for m in matches if m.get('slideIndex') == slide_index), None)

        # 获取讲稿文本
        if match and match.get('segments'):
            segments = match['segments']
        elif slide.get('note'):
            segments = split_speech_segments(slide.get('note'))
        elif slide.get('text'):
            segments = split_speech_segments(slide.get('text'))
        else:
            segments = [f"第 {slide_index} 页"]

        # 获取热区
        hotspots = slide.get('hotspots', [])
        visual_hotspots = [h for h in hotspots if h.get('kind') != 'visual']

        actions = []
        for seg_idx, segment in enumerate(segments):
            # 选择热区
            hotspot = None
            if visual_hotspots:
                if seg_idx < len(visual_hotspots):
                    hotspot = visual_hotspots[seg_idx]
                else:
                    hotspot = visual_hotspots[seg_idx % len(visual_hotspots)]

            # 添加聚光灯效果
            if hotspot and hotspot.get('kind') != 'full-slide':
                x = hotspot.get('x', 0)
                y = hotspot.get('y', 0)
                w = hotspot.get('w', 20)
                h = hotspot.get('h', 10)

                actions.append({
                    'id': f'action_{nanoid()}',
                    'type': 'spotlight',
                    'elementId': hotspot.get('id', f'hotspot_{slide_index}_{seg_idx}'),
                    'dimOpacity': 0.58
                })

                actions.append({
                    'id': f'action_{nanoid()}',
                    'type': 'laser',
                    'elementId': hotspot.get('id', f'hotspot_{slide_index}_{seg_idx}'),
                    'color': '#ff3b30'
                })

            # 添加语音动作
            actions.append({
                'id': f'action_{nanoid()}',
                'type': 'speech',
                'text': segment
            })

        action_plans.append({
            'slideIndex': slide_index,
            'actions': actions
        })

    return action_plans


def shapes_to_hotspots(shapes: List[Dict], slide_num: int) -> List[Dict]:
    hotspots = []
    for i, shape in enumerate(shapes):
        if not shape.get('text', '').strip():
            continue
        hotspots.append({
            'id': f'hotspot_s{slide_num}_{i + 1}',
            'kind': 'text',
            'text': shape['text'][:500],
            'x': shape.get('x', 0),
            'y': shape.get('y', 0),
            'w': max(shape.get('w', 10), 5),
            'h': max(shape.get('h', 5), 3),
            'priority': 90
        })

    if not hotspots:
        hotspots.append({
            'id': f'hotspot_s{slide_num}_full',
            'kind': 'full-slide',
            'text': '',
            'x': 0, 'y': 0, 'w': 100, 'h': 100,
            'priority': 0
        })
    return hotspots


# ============== ComfyUI 节点 ==============

@register_node
class OpenMAIC_PPTX导入独立版:
    """
    PPTX导入独立版 - 完全本地处理，不依赖外部服务

    功能：
    1. 将 PPTX 转换为 PDF 并提取文本
    2. 本地讲稿匹配（不需要 AI）
    3. 生成讲解动作（spotlight + laser + speech）

    类别：OpenMAIC/导入
    """
    CATEGORY = "OpenMAIC/导入"
    DISPLAY_NAME = "📊 PPTX导入（独立版）"
    RETURN_TYPES = ("DICT", "LIST", "LIST", "INT")
    RETURN_NAMES = ("课件数据", "匹配结果", "动作列表", "页面数量")
    FUNCTION = "process"

    @classmethod
    def INPUT_TYPES(cls) -> Dict[str, Any]:
        return {
            "required": {
                "PPTX路径": ("STRING", {"tooltip": "PPTX 文件的完整路径"}),
            },
            "optional": {
                "讲稿": ("STRING", {
                    "default": "",
                    "multiline": True,
                    "tooltip": "讲解稿脚本（支持标记：第1页、---分隔符）"
                }),
                "课件标题": ("STRING", {"default": "", "tooltip": "课件标题（留空用文件名）"}),
                "输出模式": (["full", "text_only"], {
                    "default": "full",
                    "tooltip": "full=含图片, text_only=仅文本"
                }),
            },
        }

    def process(
        self,
        PPTX路径: str,
        讲稿: str = "",
        课件标题: str = "",
        输出模式: str = "full",
    ) -> Tuple[Dict[str, Any], List, List, int]:
        if not PPTX路径 or not os.path.exists(PPTX路径):
            raise ValueError(f"文件不存在: {PPTX路径}")

        temp_dir = tempfile.mkdtemp(prefix="openmaic_pptx_")

        try:
            pptx_path = os.path.join(temp_dir, os.path.basename(PPTX路径))
            shutil.copy(PPTX路径, pptx_path)

            # 1. 提取 PPTX 文本
            pptx_slides = extract_pptx_text(pptx_path)
            if not pptx_slides:
                pptx_slides = extract_pptx_text_fallback(pptx_path)

            if not pptx_slides:
                raise Exception("无法从 PPTX 中提取文本内容")

            # 2. 转换 PDF
            pdf_path = convert_pptx_to_pdf(pptx_path, temp_dir)
            image_paths = []

            # 3. 构建幻灯片数据
            slides = []
            if 输出模式 == "full" and os.path.exists(pdf_path):
                image_paths = parse_pdf_to_images(pdf_path, temp_dir)
                pdf_slides = parse_pdf_extract_text(pdf_path)

                for i, slide_data in enumerate(pptx_slides):
                    slide_num = i + 1

                    # 图片
                    image_data = ""
                    if i < len(image_paths) and os.path.exists(image_paths[i]):
                        with open(image_paths[i], 'rb') as f:
                            img_base64 = base64.b64encode(f.read()).decode('utf-8')
                            image_data = f"data:image/png;base64,{img_base64}"

                    pdf_info = pdf_slides[i] if i < len(pdf_slides) else {'shapes': [], 'width': 1920, 'height': 1080}

                    text_parts = []
                    if slide_data.get('text'):
                        text_parts.append(slide_data['text'])
                    if pdf_info.get('text'):
                        text_parts.append(pdf_info['text'])

                    hotspots = shapes_to_hotspots(pdf_info.get('shapes', []), slide_num)

                    slides.append({
                        'index': slide_num,
                        'image': image_data,
                        'width': int(pdf_info.get('width', 1920)),
                        'height': int(pdf_info.get('height', 1080)),
                        'text': ' '.join(text_parts),
                        'note': slide_data.get('note', ''),
                        'hotspots': hotspots,
                    })
            else:
                for i, slide_data in enumerate(pptx_slides):
                    slide_num = i + 1
                    shapes = []
                    for shape in slide_data.get('shapes', []):
                        sw = shape.get('slide_width', 13.33)
                        sh = shape.get('slide_height', 7.5)
                        shapes.append({
                            'text': shape.get('text', ''),
                            'x': shape.get('left', 0) / sw * 100,
                            'y': shape.get('top', 0) / sh * 100,
                            'w': shape.get('width', 10) / sw * 100,
                            'h': shape.get('height', 2) / sh * 100,
                        })

                    hotspots = shapes_to_hotspots(shapes, slide_num)

                    slides.append({
                        'index': slide_num,
                        'image': '',
                        'width': 1920,
                        'height': 1080,
                        'text': slide_data.get('text', ''),
                        'note': slide_data.get('note', ''),
                        'hotspots': hotspots,
                    })

            # 4. 构建课件数据
            title = 课件标题.strip() if 课件标题 else os.path.splitext(os.path.basename(PPTX路径))[0]
            deck = {
                'title': title,
                'source_type': 'pptx',
                'slides': slides,
            }

            # 5. 讲稿匹配
            matches = match_script_to_slides(deck, 讲稿)

            # 6. 生成动作
            action_plans = build_imported_lecture_actions(deck, matches)

            # 展平动作列表
            all_actions = []
            for plan in action_plans:
                for action in plan['actions']:
                    all_actions.append(action)

            return (deck, matches, all_actions, len(slides))

        finally:
            try:
                shutil.rmtree(temp_dir)
            except Exception:
                pass


@register_node
class OpenMAIC_图片导入独立版:
    """
    图片导入独立版 - 将图片目录转为课件，本地生成动作

    类别：OpenMAIC/导入
    """
    CATEGORY = "OpenMAIC/导入"
    DISPLAY_NAME = "📁 图片导入（独立版）"
    RETURN_TYPES = ("DICT", "LIST", "LIST", "INT")
    RETURN_NAMES = ("课件数据", "匹配结果", "动作列表", "页面数量")
    FUNCTION = "process"

    @classmethod
    def INPUT_TYPES(cls) -> Dict[str, Any]:
        return {
            "required": {
                "图片目录": ("STRING", {"tooltip": "包含幻灯片图片的目录"}),
            },
            "optional": {
                "讲稿": ("STRING", {
                    "default": "",
                    "multiline": True,
                    "tooltip": "讲解稿脚本"
                }),
                "课件标题": ("STRING", {"default": "导入的课件", "tooltip": "课件标题"}),
            },
        }

    def process(
        self,
        图片目录: str,
        讲稿: str = "",
        课件标题: str = "导入的课件",
    ) -> Tuple[Dict[str, Any], List, List, int]:
        if not os.path.isdir(图片目录):
            raise ValueError(f"目录不存在: {图片目录}")

        supported_exts = {'.png', '.jpg', '.jpeg', '.webp', '.bmp'}
        image_files = [f for f in os.listdir(图片目录)
                      if os.path.splitext(f)[1].lower() in supported_exts]

        def natural_key(s):
            return [int(part) if part.isdigit() else part.lower()
                    for part in re.split(r'(\d+)', s)]

        image_files.sort(key=natural_key)

        if not image_files:
            raise ValueError(f"目录中没有找到图片文件: {图片目录}")

        slides = []
        for i, filename in enumerate(image_files):
            filepath = os.path.join(图片目录, filename)
            with open(filepath, 'rb') as f:
                img_bytes = f.read()

            ext = os.path.splitext(filename)[1].lower()
            mime = "image/jpeg" if ext in ['.jpg', '.jpeg'] else f"image/{ext.lstrip('.')}"

            try:
                from PIL import Image
                with Image.open(filepath) as img:
                    width, height = img.size
            except Exception:
                try:
                    import fitz
                    doc = fitz.open(filepath)
                    width, height = doc[0].rect.width, doc[0].rect.height
                    doc.close()
                except Exception:
                    width, height = 1920, 1080

            img_base64 = base64.b64encode(img_bytes).decode('utf-8')

            slides.append({
                'index': i + 1,
                'image': f"data:{mime};base64,{img_base64}",
                'width': int(width),
                'height': int(height),
                'text': '',
                'note': '',
                'hotspots': [{
                    'id': f'hotspot_s{i + 1}_full',
                    'kind': 'full-slide',
                    'text': '',
                    'x': 0, 'y': 0, 'w': 100, 'h': 100,
                    'priority': 0
                }],
            })

        deck = {
            'title': 课件标题,
            'source_type': 'images',
            'slides': slides,
        }

        matches = match_script_to_slides(deck, 讲稿)
        action_plans = build_imported_lecture_actions(deck, matches)
        all_actions = [action for plan in action_plans for action in plan['actions']]

        return (deck, matches, all_actions, len(slides))


# 导出
NODE_CLASS_MAPPINGS = {
    "OpenMAIC_PPTX导入独立版": OpenMAIC_PPTX导入独立版,
    "OpenMAIC_导入课件独立版": OpenMAIC_PPTX导入独立版,
    "OpenMAIC_图片导入独立版": OpenMAIC_图片导入独立版,
}

NODE_DISPLAY_NAME_MAPPINGS = {
    "OpenMAIC_PPTX导入独立版": "📊 PPTX导入（独立版）",
    "OpenMAIC_导入课件独立版": "📊 PPTX导入（独立版）",
    "OpenMAIC_图片导入独立版": "📁 图片导入（独立版）",
}
